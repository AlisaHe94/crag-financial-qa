"""Add a 'What Changes from Baseline?' comparison-table slide to the deck.

Inserts a new slide at position 6 (after Technical Contributions, before
Implementation). Updates footer page numbers from N/13 to N/14 throughout
so the deck stays internally consistent.

Run with:
    python3 _pptx_add_comparison.py [INPUT.pptx] [OUTPUT.pptx]
"""

from __future__ import annotations
import sys
import re
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# Deck palette (must match _styles.py / hero_header)
NAVY = RGBColor(0x1E, 0x27, 0x61)      # primary
TEAL = RGBColor(0x2A, 0x9D, 0x8F)      # accent (CRAG column)
SLATE_DARK = RGBColor(0x47, 0x55, 0x69) # body text
SLATE_LIGHT = RGBColor(0x94, 0xA3, 0xB8) # caption / muted
DIVIDER = RGBColor(0xE2, 0xE8, 0xF0)   # subtle line / row separator
CORAL = RGBColor(0xE7, 0x6F, 0x51)     # accent for X marks
ROW_ALT = RGBColor(0xF8, 0xFA, 0xFC)   # alternating row background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# Comparison-table content. Seven rows is the sweet spot — enough to make
# the architectural distinction concrete without overflowing one slide.
COMPARISON_ROWS = [
    ("Retrieval",
     "Single flat-text FAISS\nTop-K cosine",
     "Hybrid: separate text + table FAISS\nPer-query-type composition"),
    ("Quality gate",
     "None — top-K is final",
     "Probabilistic evaluator\nP(Relevant | q, d) ∈ [0, 1]"),
    ("Routing decision",
     "N/A — always generate",
     "3-way soft routing\nCORRECT / AMBIGUOUS / INCORRECT"),
    ("Query classification",
     "None",
     "LLM classifier\ntable_lookup / narrative / multimodal / OOC"),
    ("Out-of-corpus handling",
     "Hallucinate or refuse",
     "Tavily web fallback\nSourced live answer"),
    ("Multi-part synthesis",
     "Single retrieval + single generation",
     "Completeness check + sharper re-prompt\n+ optional web augmentation"),
    ("Interpretability",
     "Just an answer (opaque)",
     "Visible badges: query type, routing,\ntier, confidence + readable citations"),
]


def add_textbox(slide, left, top, width, height, text,
                font_size=12, bold=False, color=SLATE_DARK,
                align=PP_ALIGN.LEFT, font_name="Calibri",
                anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def set_cell_fill(cell, rgb: RGBColor):
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb


def set_cell_text(cell, text: str, *, font_size=11, bold=False,
                  color=SLATE_DARK, align=PP_ALIGN.LEFT, font_name="Calibri"):
    """Replace cell text and apply formatting. Preserves multi-line content
    using paragraph breaks on \n."""
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0.08 * 914400)   # ~0.08"
    tf.margin_right = Emu(0.08 * 914400)
    tf.margin_top = Emu(0.05 * 914400)
    tf.margin_bottom = Emu(0.05 * 914400)
    # Clear existing paragraphs
    for p in list(tf.paragraphs):
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
    # First line goes in the existing first paragraph
    lines = text.split("\n")
    p0 = tf.paragraphs[0]
    p0.alignment = align
    r0 = p0.add_run()
    r0.text = lines[0]
    r0.font.name = font_name
    r0.font.size = Pt(font_size)
    r0.font.bold = bold
    r0.font.color.rgb = color
    # Subsequent lines as new paragraphs
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font_name
        r.font.size = Pt(font_size)
        r.font.bold = bold
        r.font.color.rgb = color


def add_comparison_slide(prs):
    """Add a new comparison-table slide using a blank layout."""
    # The deck has a single 'DEFAULT' layout — use it.
    blank = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank)

    SLIDE_W = prs.slide_width
    SLIDE_H = prs.slide_height
    margin = Inches(0.5)

    # --- Section eyebrow / category label ---
    add_textbox(
        slide, margin, Inches(0.5),
        SLIDE_W - 2 * margin, Inches(0.3),
        "ARCHITECTURE COMPARISON",
        font_size=11, bold=True, color=SLATE_LIGHT,
        font_name="Calibri",
    )

    # --- Title ---
    add_textbox(
        slide, margin, Inches(0.85),
        SLIDE_W - 2 * margin, Inches(0.7),
        "What Changes from Baseline?",
        font_size=32, bold=True, color=NAVY,
        font_name="Georgia",
    )

    # --- Subtitle ---
    add_textbox(
        slide, margin, Inches(1.6),
        SLIDE_W - 2 * margin, Inches(0.4),
        "Component-by-component diff. Same corpus, same generator — only the architecture differs.",
        font_size=14, color=SLATE_DARK, font_name="Calibri",
    )

    # --- Comparison table ---
    rows = len(COMPARISON_ROWS) + 1  # +1 for header
    cols = 3
    table_left = margin
    table_top = Inches(2.2)
    table_width = SLIDE_W - 2 * margin
    table_height = Inches(4.6)

    table_shape = slide.shapes.add_table(
        rows, cols, table_left, table_top, table_width, table_height
    )
    table = table_shape.table

    # Column widths (component | baseline | crag)
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(4.4)
    # Compute remaining width in EMU then assign
    remaining = int(table_width) - int(Inches(2.2)) - int(Inches(4.4))
    table.columns[2].width = remaining

    # --- Header row ---
    headers = ["Component", "Baseline RAG", "Probabilistic CRAG (Ours)"]
    header_colors = [NAVY, NAVY, TEAL]  # CRAG column distinct
    for c, (text, fill) in enumerate(zip(headers, header_colors)):
        cell = table.cell(0, c)
        set_cell_fill(cell, fill)
        set_cell_text(
            cell, text,
            font_size=12, bold=True, color=WHITE,
            align=PP_ALIGN.LEFT, font_name="Calibri",
        )
    # Slight extra height for header
    table.rows[0].height = Inches(0.45)

    # --- Data rows ---
    for i, (component, baseline, crag) in enumerate(COMPARISON_ROWS, start=1):
        # Alternating row backgrounds for readability
        bg = ROW_ALT if i % 2 == 0 else WHITE
        for c in range(3):
            cell = table.cell(i, c)
            set_cell_fill(cell, bg)
        # Cell contents
        set_cell_text(
            table.cell(i, 0), component,
            font_size=11, bold=True, color=NAVY,
        )
        set_cell_text(
            table.cell(i, 1), baseline,
            font_size=10, color=SLATE_DARK,
        )
        set_cell_text(
            table.cell(i, 2), crag,
            font_size=10, color=SLATE_DARK,
        )
        table.rows[i].height = Inches(0.55)

    # --- Footer (matches deck convention: "STAT 5293 GenAI · Probabilistic CRAG · X/N") ---
    add_textbox(
        slide,
        margin,
        SLIDE_H - Inches(0.45),
        SLIDE_W - 2 * margin,
        Inches(0.3),
        "STAT 5293 GenAI · Probabilistic CRAG · 6/14",
        font_size=9, color=SLATE_LIGHT,
        align=PP_ALIGN.LEFT, font_name="Calibri",
    )

    return slide


def move_slide(prs, src_index: int, dst_index: int):
    """Reorder slides via the underlying XML — python-pptx has no public API."""
    sldIdLst = prs.slides._sldIdLst
    children = list(sldIdLst)
    moving = children[src_index]
    sldIdLst.remove(moving)
    # If we removed something before dst, dst doesn't shift because we list-copied first
    sldIdLst.insert(dst_index, moving)


# Footer line pattern — matches "STAT 5293 GenAI · Probabilistic CRAG · 5/13"
FOOTER_PATTERN = re.compile(
    r"(STAT\s*5293\s*GenAI\s*·\s*Probabilistic\s*CRAG\s*·\s*)(\d+)\s*/\s*13"
)


def renumber_footers(prs, new_total: int = 14, new_slide_position: int = 6):
    """Walk every text run and renumber 'X/13' footers to reflect 14 slides.

    Slides at user-facing position < new_slide_position keep their number.
    Slides at user-facing position >= new_slide_position get +1 because the
    new comparison slide was inserted before them.
    """
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    m = FOOTER_PATTERN.search(run.text)
                    if not m:
                        continue
                    old_num = int(m.group(2))
                    # The footer's claimed number reflects the slide's position
                    # before insertion. Now its actual position is slide_idx.
                    # Just rewrite to "<slide_idx>/<new_total>".
                    new_text = f"{m.group(1)}{slide_idx}/{new_total}"
                    run.text = FOOTER_PATTERN.sub(new_text, run.text)


def main() -> None:
    in_path = sys.argv[1] if len(sys.argv) > 1 else "Probabilistic_CRAG_Final_Presentation_v2.pptx"
    out_path = sys.argv[2] if len(sys.argv) > 2 else "Probabilistic_CRAG_Final_Presentation_v3.pptx"

    prs = Presentation(in_path)
    n_before = len(prs.slides)
    print(f"Loaded {in_path}: {n_before} slides")

    add_comparison_slide(prs)
    print(f"Added comparison slide at end (position {len(prs.slides)})")

    # Move from end (index n_before, since we just appended) to position 5 (0-indexed),
    # which becomes user-facing slide 6
    move_slide(prs, src_index=n_before, dst_index=5)
    print("Moved comparison slide to user-facing position 6")

    # Renumber footers from X/13 to N/14 based on new positions
    renumber_footers(prs, new_total=14, new_slide_position=6)
    print("Renumbered footers to N/14")

    prs.save(out_path)
    print(f"Saved -> {out_path}")


if __name__ == "__main__":
    main()
