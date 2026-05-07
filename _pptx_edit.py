"""One-shot edit script for Probabilistic_CRAG_Final_Presentation.pptx.

Run with:
    python3 _pptx_edit.py [INPUT.pptx] [OUTPUT.pptx]

Edits applied (preserving slide layouts, fonts, colors):
  - Slide 5  : "Markdown chunks" -> "key-value sentence chunks"
  - Slide 8  : tau_high 0.55 -> 0.75; fix strong-corpus-match guard wording
  - Slide 9  : replace demo walkthrough with v3 lineup
  - Slide 12 : replace MS-MARCO MiniLM mention; move dynamic-routing
               from future work to implemented; add strip-level refinement
"""

from __future__ import annotations
import sys
import re
from copy import deepcopy
from pptx import Presentation
from pptx.util import Pt


def replace_text_in_run(run, old: str, new: str) -> bool:
    """Replace text within a single run while preserving formatting."""
    if old in run.text:
        run.text = run.text.replace(old, new)
        return True
    return False


def replace_in_shape(shape, replacements: list[tuple[str, str]]) -> int:
    """Walk all text frames in a shape and apply replacements."""
    n = 0
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for old, new in replacements:
                    if replace_text_in_run(run, old, new):
                        n += 1
    if shape.shape_type == 19:  # table
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            for old, new in replacements:
                                if replace_text_in_run(run, old, new):
                                    n += 1
    return n


def replace_paragraph_text(paragraph, new_text: str) -> None:
    """Replace the entire text of a paragraph, preserving formatting of the
    first run. Subsequent runs are cleared."""
    if not paragraph.runs:
        # Create a run if none exist
        from pptx.oxml.ns import qn
        r = paragraph._p.makeelement(qn("a:r"), {})
        paragraph._p.append(r)
        return
    first = paragraph.runs[0]
    first.text = new_text
    # Remove subsequent runs
    for run in paragraph.runs[1:]:
        run._r.getparent().remove(run._r)


def find_paragraph_by_prefix(text_frame, prefix: str):
    """Return the paragraph whose text starts with prefix (case-sensitive),
    else None."""
    for para in text_frame.paragraphs:
        if para.text.strip().startswith(prefix):
            return para
    return None


def find_paragraph_containing(text_frame, needle: str):
    for para in text_frame.paragraphs:
        if needle in para.text:
            return para
    return None


# ---------------------------------------------------------------------------
# Per-slide edits
# ---------------------------------------------------------------------------

def edit_slide5(slide) -> None:
    """Slide 5: 'Markdown chunks' -> 'key-value sentence chunks'."""
    for shape in slide.shapes:
        replace_in_shape(shape, [
            ("Markdown chunks", "key-value sentence chunks"),
            ("Markdown", "key-value sentence"),
        ])


def edit_slide8(slide) -> None:
    """Slide 8:
      - tau_high (0.55) -> tau_high (0.75)
      - Fix 'score < tau_low OR max-cosine < 0.65' wording
        The actual logic: web fires when (low confidence) AND (max-cosine < 0.65).
        max-cosine >= 0.65 is a strong-corpus-match GUARD that prevents web fallback.
    """
    for shape in slide.shapes:
        replace_in_shape(shape, [
            ("τ_high (0.55)", "τ_high (0.75)"),
            ("score < τ_low  OR  max-cosine < 0.65",
             "score < τ_low  AND  max-cosine < 0.65"),
            ("Incorrect: replace context with Tavily web search.",
             "Web fallback fires only when corpus is weak. "
             "max-cosine ≥ 0.65 is a strong-corpus-match guard that "
             "keeps CRAG on corpus when at least one chunk clearly fits."),
        ])


def edit_slide9(slide) -> None:
    """Slide 9: replace the demo walkthrough with the v3 lineup.

    Original three rows (TYPE-3 / TYPE-1 / OUT-OF-CORPUS) are replaced with:
      MULTIMODAL · Meta DAU + Reality Labs        — interpretability win
      PARTIAL OOC · Microsoft Azure + stock price — completeness check + web aug
      OUT-OF-CORPUS · Bitcoin price                — pure web fallback
    """
    # Map of old-line -> new-line replacements (run-level so we keep formatting).
    # We target the distinctive substrings in each row's title / CRAG / Baseline lines.
    replacements: list[tuple[str, str]] = [
        # Row labels
        ("TYPE-3 · TABLE", "TYPE-4 · MULTIMODAL"),
        ("TYPE-1 · TEXT", "PARTIAL OUT-OF-CORPUS"),
        ("OUT-OF-CORPUS", "OUT-OF-CORPUS"),

        # Question texts (replace whole strings to avoid partial matches)
        ("“What were Apple's total net sales in FY 2025?”",
         "“How did Meta's family of apps DAU change in 2025, and what does "
         "management say about Reality Labs spending?”"),
        ("“How does Microsoft describe its Azure business in the FY 2025 10-K?”",
         "“What was Microsoft's Azure revenue growth in FY 2025, and what is "
         "Microsoft's current stock price?”"),
        ("“What is the current federal funds rate?”",
         "“What is the current price of Bitcoin?”"),

        # CRAG row body lines
        ("✓ $416,161M · routes to AMBIGUOUS · pulls from extracted income-statement table chunk",
         "✓ Tier text+table · routing CORRECT · DAU 3.58B from table chunk + Reality Labs prose synthesized"),
        ("✓ Detailed: Server products & cloud services revenue $98.4B, +23% growth driven by Azure +34%",
         "✓ Tier text+table+web · completeness check fires · Azure 34% growth from corpus + live stock price from web"),
        ("✓ INCORRECT routing → Tavily fallback → \"3.50–3.75% (Trading Economics, NerdWallet)\"",
         "✓ Tier web (OOC) · classified out_of_corpus · routes to Tavily · sourced live price"),

        # Baseline row body lines
        ("✓ $416,161M · top-K text chunk happens to contain the figure flat-text",
         "≈ Both halves answered at top_k=5 · but no routing badge, no tier, no confidence — opaque"),
        ("✗ \"Unable to verify\" — top-4 text chunks miss the segment-revenue paragraph",
         "✗ Cannot reach the web — refuses or hallucinates the stock price from training data"),
        ("✗ \"I cannot verify\" — no fallback path exists in baseline RAG",
         "✗ No fallback path — refuses or hallucinates"),
    ]
    for shape in slide.shapes:
        replace_in_shape(shape, replacements)


def edit_slide12(slide) -> None:
    """Slide 12: limitations and future work.

    - Cross-encoder limitation now mentions BAAI/bge-reranker-base
      and uncalibrated scores (rather than MS-MARCO MiniLM).
    - "Dynamic LLM-classified routing" moves from Next Steps to
      implemented (we ship it). Replaced with strip-level refinement.
    """
    replacements: list[tuple[str, str]] = [
        # Cross-encoder ceiling — we replaced MS-MARCO MiniLM with bge-reranker-base.
        # Update both the heading-line description and the body sentence.
        ("MS-MARCO MiniLM was trained on web search, not financial filings — discrimination plateaus.",
         "BAAI/bge-reranker-base scores are uncalibrated on financial language; "
         "routing relies on multiple guard rails (cosine ordering, strong-match guard) "
         "to compensate. Calibrating on a labeled validation set is open work."),

        # Dynamic LLM routing -> implemented; strip-level refinement -> new future work
        ("Dynamic LLM-classified routing",
         "Strip-level knowledge refinement"),
        ("Use the LLM to pre-classify each query as text / table / OOC and bias retrieval accordingly.",
         "Score sentence-level strips inside chunks (the original CRAG paper's contribution we did not yet implement) — directly addresses lost-in-the-middle failures on long MD&A passages."),

        # Threshold calibration — refresh the body sentence to be current
        ("Mean confidence sits ~0.55–0.60 across most queries; routing badge is too often AMBIGUOUS to be informative.",
         "Mean confidence is dragged down by reranker noise; routing badge often "
         "shows AMBIGUOUS even when one chunk strongly matches. A calibrated "
         "evaluator (logistic regression on labeled (q, d) pairs) would give "
         "actual probabilities rather than hand-tuned weighted sums."),
    ]
    for shape in slide.shapes:
        replace_in_shape(shape, replacements)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> None:
    in_path = sys.argv[1] if len(sys.argv) > 1 else "Probabilistic_CRAG_Final_Presentation.pptx"
    out_path = sys.argv[2] if len(sys.argv) > 2 else "Probabilistic_CRAG_Final_Presentation_v2.pptx"

    pres = Presentation(in_path)
    slides = list(pres.slides)
    if len(slides) < 13:
        print(f"WARNING: deck has {len(slides)} slides, expected 13", file=sys.stderr)

    # Slides are 0-indexed in python-pptx; user-facing is 1-indexed
    edit_slide5(slides[4])
    edit_slide8(slides[7])
    edit_slide9(slides[8])
    edit_slide12(slides[11])

    pres.save(out_path)
    print(f"Saved updated deck to: {out_path}")


if __name__ == "__main__":
    main()
